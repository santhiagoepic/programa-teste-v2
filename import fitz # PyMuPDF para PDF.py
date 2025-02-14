import fitz  # PyMuPDF para PDF
import pandas as pd  # Pandas para Excel
import os
import docx  # python-docx para Word
import pptx  # python-pptx para PowerPoint
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from ttkthemes import ThemedTk
from tqdm import tqdm  # Para barra de progresso no terminal
import threading  # Para rodar a busca em uma thread separada

# Variáveis globais para controle da busca
search_in_progress = False
cancel_search = False

# Função para buscar a palavra-chave nos arquivos selecionados
def search_files(directory, keyword, search_pdf, search_docx, search_xlsx, search_pptx, progress_var, progress_label, result_text, show_file_info):
    global search_in_progress, cancel_search
    search_in_progress = True
    cancel_search = False
    results = []
    files = []

    for root, _, filenames in os.walk(directory):
        for file in filenames:
            if (search_pdf and file.endswith(".pdf")) or \
               (search_docx and file.endswith(".docx")) or \
               (search_xlsx and file.endswith(".xlsx")) or \
               (search_pptx and file.endswith(".pptx")):
                files.append(os.path.join(root, file))

    total_files = len(files)
    progress_step = 100 / total_files if total_files else 0
    progress = 0

    for file_path in tqdm(files, desc="Searching", unit="file"):
        if cancel_search:
            break

        file_name = os.path.basename(file_path)
        found = False

        file_info = ""
        if show_file_info.get():
            mod_time = os.path.getmtime(file_path)
            mod_time_str = pd.to_datetime(mod_time, unit='s').strftime('%Y-%m-%d %H:%M:%S')
            file_info = f" | Modificado: {mod_time_str} | Local: {file_path}"

        if keyword.lower() in file_name.lower():
            results.append(f"{file_name} (Nome do arquivo){file_info}")
            found = True

        if file_path.endswith(".pdf") and search_pdf:
            try:
                with fitz.open(file_path) as pdf_doc:
                    for page_num in range(len(pdf_doc)):
                        text = pdf_doc[page_num].get_text()
                        if keyword.lower() in text.lower():
                            results.append(f"{file_name} - Página {page_num + 1}{file_info}")
                            found = True
                            break
            except Exception as e:
                print(f"Erro ao processar {file_path}: {e}")

        elif file_path.endswith(".docx") and search_docx:
            try:
                doc = docx.Document(file_path)
                for para in doc.paragraphs:
                    if keyword.lower() in para.text.lower():
                        results.append(f"{file_name} (Word){file_info}")
                        found = True
                        break
            except Exception as e:
                print(f"Erro ao processar {file_path}: {e}")

        elif file_path.endswith(".xlsx") and search_xlsx:
            try:
                df = pd.read_excel(file_path, engine='openpyxl')
                if df.astype(str).apply(lambda x: x.str.contains(keyword, case=False, na=False)).any().any():
                    results.append(f"{file_name} (Excel){file_info}")
                    found = True
            except Exception as e:
                print(f"Erro ao ler {file_path}: {e}")

        elif file_path.endswith(".pptx") and search_pptx:
            try:
                presentation = pptx.Presentation(file_path)
                for slide_num, slide in enumerate(presentation.slides, start=1):
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and keyword.lower() in shape.text.lower():
                            results.append(f"{file_name} - Slide {slide_num}{file_info}")
                            found = True
                            break
                    if found:
                        break
            except Exception as e:
                print(f"Erro ao processar {file_path}: {e}")

        progress += progress_step
        progress_var.set(progress)
        progress_label.update()

    search_in_progress = False
    result_text.delete(1.0, tk.END)
    result_text.insert(tk.END, "\n".join(results) if results else "Nenhum resultado encontrado.")

# Função para iniciar a busca em uma thread separada
def start_search():
    if search_in_progress:
        messagebox.showwarning("Aviso", "A busca já está em andamento.")
        return
    
    directory = directory_entry.get()
    keyword = keyword_entry.get()
    if not directory or not keyword:
        messagebox.showerror("Erro", "Por favor, insira um diretório e uma palavra-chave.")
        return

    # Iniciar a busca em uma thread separada
    threading.Thread(target=search_files, args=(directory, keyword, search_pdf.get(), search_docx.get(), search_xlsx.get(), search_pptx.get(), progress_var, progress_label, result_text, show_file_info)).start()

# Função para cancelar a busca
def cancel_search_func():
    global cancel_search
    if search_in_progress:
        cancel_search = True
        messagebox.showinfo("Busca Cancelada", "A busca foi cancelada.")
    else:
        messagebox.showinfo("Busca Não Iniciada", "Nenhuma busca foi iniciada.")

# Interface gráfica
root = ThemedTk(theme="arc")
root.title("Busca de Arquivos")
root.geometry("720x720")

frame = ttk.Frame(root, padding=10)
frame.pack(fill=tk.BOTH, expand=True)

keyword_label = ttk.Label(frame, text="Palavra-chave:", width=20, anchor="w")
keyword_label.pack()
keyword_entry = ttk.Entry(frame)
keyword_entry.pack(fill=tk.X)

directory_label = ttk.Label(frame, text="Diretório:", width=20, anchor="w")
directory_label.pack()
directory_entry = ttk.Entry(frame)
directory_entry.pack(fill=tk.X)

def browse_directory():
    directory = filedialog.askdirectory()
    if directory:
        directory_entry.delete(0, tk.END)
        directory_entry.insert(0, directory)

browse_button = ttk.Button(frame, text="Selecionar Pasta", command=browse_directory)
browse_button.pack()

search_pdf = tk.BooleanVar()
search_docx = tk.BooleanVar()
search_xlsx = tk.BooleanVar()
search_pptx = tk.BooleanVar()

pdf_check = ttk.Checkbutton(frame, text="PDF", variable=search_pdf)
pdf_check.pack()
docx_check = ttk.Checkbutton(frame, text="Word", variable=search_docx)
docx_check.pack()
xlsx_check = ttk.Checkbutton(frame, text="Excel", variable=search_xlsx)
xlsx_check.pack()
pptx_check = ttk.Checkbutton(frame, text="PowerPoint", variable=search_pptx)
pptx_check.pack()

show_file_info = tk.BooleanVar()

info_check = ttk.Checkbutton(frame, text="Mostrar informações do arquivo", variable=show_file_info)
info_check.pack()

progress_var = tk.DoubleVar()
progress_bar = ttk.Progressbar(frame, variable=progress_var, maximum=100)
progress_bar.pack(fill=tk.X, pady=5)
progress_label = ttk.Label(frame, text="Progresso")
progress_label.pack()

result_text = tk.Text(frame, height=10)
result_text.pack(fill=tk.BOTH, expand=True)

search_button = ttk.Button(frame, text="Buscar", command=start_search)
search_button.pack()

cancel_button = ttk.Button(frame, text="Cancelar Busca", command=cancel_search_func)
cancel_button.pack()

root.mainloop()